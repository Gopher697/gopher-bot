"""
Discord bridge for Gopher-bot.

Receives messages from a Discord channel and routes them through Awareness.
Modelled on the GopherVault discord_vault_bridge.py structure.

Setup
-----
1. Create a Discord application + bot at https://discord.com/developers/applications
2. Enable "Message Content Intent" under Bot → Privileged Gateway Intents
3. Add  DISCORD_BOT_TOKEN = "<your token>"  to world_models/config.py
4. Optionally add  DISCORD_CHANNEL = "your-channel-name"  (default: "gopher")
5. Invite the bot to your server with scopes: bot, permissions: Send Messages / Read Message History
6. Run:  python interface/discord_bot.py

The bridge also starts automatically from server.py if a token is configured
and the module is importable (no extra step needed when running the full stack).
"""

from __future__ import annotations

import asyncio
import sys
from datetime import datetime, timezone
from pathlib import Path

try:
    import discord
except ModuleNotFoundError:
    print("Missing package: discord.py")
    print("Install with:  pip install discord.py")
    raise SystemExit(1)

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from interface import bot  # noqa: E402


# ── Config helpers ─────────────────────────────────────────────────────────


def _config_str(attr: str, env_var: str, default: str = "") -> str:
    """Read a string setting from world_models/config.py, then env, then default."""
    import os

    try:
        from world_models import config

        value = str(getattr(config, attr, "")).strip()
        if value:
            return value
    except Exception:
        pass
    return os.environ.get(env_var, default).strip()


DISCORD_TOKEN: str = _config_str("DISCORD_BOT_TOKEN", "GOPHER_DISCORD_TOKEN")
DISCORD_CHANNEL: str = _config_str("DISCORD_CHANNEL", "GOPHER_DISCORD_CHANNEL", "gopher")

MAX_REPLY_CHARS = 1900          # Discord hard limit is 2000; leave headroom
MAX_ATTACHMENT_BYTES = 25 * 1024 * 1024
RATE_LIMIT_MAX_MESSAGES = 30
RATE_LIMIT_WINDOW_SECONDS = 60
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".heic", ".heif", ".bmp"}
AUDIO_EXTENSIONS = {".ogg", ".mp3", ".wav", ".m4a", ".flac", ".webm", ".opus"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".webm", ".avi", ".mkv", ".wmv"}
WEBM_AUDIO_MAX_BYTES = 10 * 1024 * 1024
PROACTIVE_CHECK_INTERVAL = 10          # seconds between polls
PROACTIVE_RATE_LIMIT_SECONDS = 60      # minimum gap between proactive Discord sends


# ── Internal state ─────────────────────────────────────────────────────────

_rate_state: dict[str, dict] = {}
_process_lock: asyncio.Lock | None = None   # created inside the event loop
_last_proactive_at: float = 0.0
_last_proactive_message_id: int = 0


def _get_process_lock() -> asyncio.Lock:
    global _process_lock
    if _process_lock is None:
        _process_lock = asyncio.Lock()
    return _process_lock


async def _proactive_loop() -> None:
    """
    Poll server.py's proactive message endpoint and forward new messages.

    BrainLoop already handles reminder firing and bid draining in server.py.
    This bridge only polls GET /proactive-messages?since=<id>; it does not
    touch the bid queue or Neo4j.
    """
    import json as _json
    import time as _time
    import urllib.request as _urllib_request

    global _last_proactive_at, _last_proactive_message_id

    await client.wait_until_ready()

    target_channel = None
    for ch in client.get_all_channels():
        if getattr(ch, "name", "") == DISCORD_CHANNEL:
            target_channel = ch
            break

    if target_channel is None:
        print(
            f"[discord] Proactive loop: channel #{DISCORD_CHANNEL} not found -- "
            "proactive messages disabled"
        )
        return

    print(f"[discord] Proactive loop started (channel: #{DISCORD_CHANNEL})")

    while not client.is_closed():
        await asyncio.sleep(PROACTIVE_CHECK_INTERVAL)
        try:
            now = _time.time()
            if now - _last_proactive_at < PROACTIVE_RATE_LIMIT_SECONDS:
                continue

            url = (
                "http://localhost:5000/proactive-messages"
                f"?since={_last_proactive_message_id}"
            )

            def _fetch() -> dict:
                with _urllib_request.urlopen(url, timeout=5) as resp:
                    return _json.loads(resp.read())

            payload = await asyncio.to_thread(_fetch)
            messages = payload.get("messages") or []

            for msg in messages:
                msg_id = int(msg.get("id") or 0)
                text = str(msg.get("text") or "").strip()
                if not text:
                    _last_proactive_message_id = max(
                        _last_proactive_message_id,
                        msg_id,
                    )
                    continue

                for chunk in _split_reply(text):
                    await target_channel.send(chunk)

                _last_proactive_message_id = msg_id
                _last_proactive_at = now
                break

        except Exception as exc:
            print(f"[discord] Proactive loop error: {exc}")


# ── Utilities ──────────────────────────────────────────────────────────────


def _check_rate_limit(user_id: int, now: datetime) -> tuple[bool, bool]:
    """Return (allowed, should_warn)."""
    key = str(user_id)
    state = _rate_state.get(key)
    if state is None or (now - state["window_start"]).total_seconds() >= RATE_LIMIT_WINDOW_SECONDS:
        _rate_state[key] = {"window_start": now, "count": 1, "warned": False}
        return True, False
    state["count"] += 1
    if state["count"] <= RATE_LIMIT_MAX_MESSAGES:
        return True, False
    if not state["warned"]:
        state["warned"] = True
        return False, True
    return False, False


def _split_reply(text: str) -> list[str]:
    """Split a long reply into Discord-sized chunks, preferring newline boundaries."""
    text = (text or "").strip()
    if not text:
        return []
    chunks: list[str] = []
    while len(text) > MAX_REPLY_CHARS:
        split_at = text.rfind("\n", 0, MAX_REPLY_CHARS)
        if split_at < 500:
            split_at = MAX_REPLY_CHARS
        chunks.append(text[:split_at].strip())
        text = text[split_at:].strip()
    if text:
        chunks.append(text)
    return chunks


def _is_audio_attachment(suffix: str, size: int | None = None) -> bool:
    if suffix not in AUDIO_EXTENSIONS:
        return False
    if suffix == ".webm" and size and size > WEBM_AUDIO_MAX_BYTES:
        return False
    return True


def _is_video_attachment(suffix: str, size: int | None = None) -> bool:
    if suffix not in VIDEO_EXTENSIONS:
        return False
    if _is_audio_attachment(suffix, size):
        return False
    return True


def _extract_document_text(data: bytes, filename: str) -> str | None:
    """
    Attempt to extract readable text from a binary document.

    Tries format-specific parsers based on file extension. Returns the extracted
    text string (possibly empty) on success, or None if the format is unsupported
    or the required library is not installed.

    Supported formats:
        .pdf            - pdfplumber
        .docx           - python-docx
        .xlsx / .xls    - openpyxl
        .pptx           - python-pptx
        .csv            - decoded as UTF-8
        .rtf            - basic tag stripping (no extra dependency)
    """
    suffix = Path(filename).suffix.lower()

    # --- PDF ---
    if suffix == ".pdf":
        try:
            import pdfplumber
            from io import BytesIO as _BytesIO

            text_parts: list[str] = []
            with pdfplumber.open(_BytesIO(data)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        text_parts.append(page_text.strip())
            return "\n\n".join(text_parts) if text_parts else ""
        except ImportError:
            print("[discord] pdfplumber not installed; cannot parse PDF")
            return None
        except Exception as exc:
            print(f"[discord] PDF parse failed for {filename}: {exc}")
            return None

    # --- Word (.docx) ---
    if suffix == ".docx":
        try:
            import docx
            from io import BytesIO as _BytesIO

            doc = docx.Document(_BytesIO(data))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs)
        except ImportError:
            print("[discord] python-docx not installed; cannot parse .docx")
            return None
        except Exception as exc:
            print(f"[discord] DOCX parse failed for {filename}: {exc}")
            return None

    # --- Excel (.xlsx / .xls) ---
    if suffix in {".xlsx", ".xls"}:
        try:
            import openpyxl
            from io import BytesIO as _BytesIO

            wb = openpyxl.load_workbook(_BytesIO(data), read_only=True, data_only=True)
            rows: list[str] = []
            for sheet in wb.worksheets:
                rows.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    if any(c.strip() for c in cells):
                        rows.append("\t".join(cells))
            return "\n".join(rows)
        except ImportError:
            print("[discord] openpyxl not installed; cannot parse .xlsx")
            return None
        except Exception as exc:
            print(f"[discord] Excel parse failed for {filename}: {exc}")
            return None

    # --- PowerPoint (.pptx) ---
    if suffix == ".pptx":
        try:
            from pptx import Presentation
            from io import BytesIO as _BytesIO

            prs = Presentation(_BytesIO(data))
            slides: list[str] = []
            for i, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides)
        except ImportError:
            print("[discord] python-pptx not installed; cannot parse .pptx")
            return None
        except Exception as exc:
            print(f"[discord] PPTX parse failed for {filename}: {exc}")
            return None

    # --- CSV fallback for non-UTF-8 retry path ---
    if suffix == ".csv":
        try:
            return data.decode("utf-8", errors="strict")
        except UnicodeDecodeError:
            return None

    # --- RTF (basic tag strip, no extra dependency) ---
    if suffix == ".rtf":
        try:
            import re as _re

            text = data.decode("latin-1", errors="replace")
            text = _re.sub(r"\\[a-z]+[-\d]*\s?", " ", text)
            text = _re.sub(r"[{}\\]", "", text)
            text = " ".join(text.split())
            return text
        except Exception as exc:
            print(f"[discord] RTF strip failed for {filename}: {exc}")
            return None

    # Unsupported binary format
    return None


async def _read_all_text_attachments(
    message: discord.Message,
) -> tuple[str, list[dict]]:
    """
    Read all non-image attachments.

    Returns:
        (combined_text, structured_list)
        combined_text: concatenated content for current-turn context.
        structured_list: [{"filename": str, "content": str}, ...] for graph ingestion.
    """
    parts: list[str] = []
    structured: list[dict] = []
    for attachment in message.attachments:
        suffix = Path(attachment.filename or "").suffix.lower()
        if (
            suffix in IMAGE_EXTENSIONS
            or _is_audio_attachment(suffix, attachment.size)
            or _is_video_attachment(suffix, attachment.size)
        ):
            continue
        if attachment.size and attachment.size > MAX_ATTACHMENT_BYTES:
            print(f"[discord] Skipping oversized attachment: {attachment.filename}")
            parts.append(f"[{attachment.filename}]: (file too large to transmit)")
            continue
        try:
            data = await attachment.read()
        except Exception as exc:
            print(f"[discord] Failed to download {attachment.filename}: {exc}")
            parts.append(f"[{attachment.filename}]: (download failed)")
            continue
        try:
            text = data.decode("utf-8", errors="strict")
            parts.append(f"[{attachment.filename}]:\n{text}")
            structured.append({"filename": attachment.filename, "content": text})
        except UnicodeDecodeError:
            extracted = _extract_document_text(data, attachment.filename or "")
            if extracted is not None:
                if extracted.strip():
                    parts.append(f"[{attachment.filename}]:\n{extracted}")
                    structured.append(
                        {"filename": attachment.filename, "content": extracted}
                    )
                else:
                    parts.append(
                        f"[{attachment.filename}]: "
                        "(document contained no extractable text)"
                    )
            else:
                parts.append(
                    f"[{attachment.filename}]: "
                    "(binary file -- format not supported for text extraction)"
                )
    return "\n\n".join(parts), structured


async def _download_image_attachments(message: discord.Message) -> list[dict]:
    """
    Download image attachments and return them as a list of
    {"filename": str, "data": bytes} dicts for the Sensory coordinator.
    Skips files over MAX_ATTACHMENT_BYTES.
    """
    result = []
    for attachment in message.attachments:
        suffix = Path(attachment.filename or "").suffix.lower()
        if suffix not in IMAGE_EXTENSIONS:
            continue
        if attachment.size and attachment.size > MAX_ATTACHMENT_BYTES:
            print(f"[discord] Skipping oversized image: {attachment.filename}")
            continue
        try:
            data = await attachment.read()
            result.append({"filename": attachment.filename, "data": data})
        except Exception as exc:
            print(f"[discord] Failed to download {attachment.filename}: {exc}")
    return result


async def _download_audio_attachments(message: discord.Message) -> list[dict]:
    """
    Download audio attachments and return them as
    {"filename": str, "data": bytes} dicts for the Sensory coordinator.
    """
    result = []
    for attachment in message.attachments:
        suffix = Path(attachment.filename or "").suffix.lower()
        if not _is_audio_attachment(suffix, attachment.size):
            continue
        if attachment.size and attachment.size > MAX_ATTACHMENT_BYTES:
            print(f"[discord] Skipping oversized audio: {attachment.filename}")
            continue
        try:
            data = await attachment.read()
            result.append({"filename": attachment.filename, "data": data})
        except Exception as exc:
            print(f"[discord] Failed to download {attachment.filename}: {exc}")
    return result


async def _download_video_attachments(message: discord.Message) -> list[dict]:
    """
    Download video attachments and return them as
    {"filename": str, "data": bytes} dicts for the Sensory coordinator.
    """
    result = []
    for attachment in message.attachments:
        suffix = Path(attachment.filename or "").suffix.lower()
        if not _is_video_attachment(suffix, attachment.size):
            continue
        if attachment.size and attachment.size > MAX_ATTACHMENT_BYTES:
            print(f"[discord] Skipping oversized video: {attachment.filename}")
            continue
        try:
            data = await attachment.read()
            result.append({"filename": attachment.filename, "data": data})
        except Exception as exc:
            print(f"[discord] Failed to download {attachment.filename}: {exc}")
    return result


# ── Discord client ─────────────────────────────────────────────────────────

intents = discord.Intents.default()
intents.message_content = True
client = discord.Client(intents=intents)


@client.event
async def on_ready() -> None:
    print(f"[discord] Gopher-bot bridge online as {client.user}")
    print(f"[discord] Listening on channel: #{DISCORD_CHANNEL}")
    asyncio.create_task(_proactive_loop())


@client.event
async def on_message(message: discord.Message) -> None:
    # Ignore bots (including ourselves)
    if message.author.bot:
        return

    # Only respond in the configured channel
    if getattr(message.channel, "name", "") != DISCORD_CHANNEL:
        return

    now = datetime.now(tz=timezone.utc)
    allowed, warn = _check_rate_limit(message.author.id, now)
    if not allowed:
        if warn:
            await message.reply("Slow down — too many messages.", mention_author=False)
        return

    async with _get_process_lock():
        try:
            # Combine message text with any readable non-image file attachments
            content = (message.content or "").strip()
            attachment_text, text_attachments = await _read_all_text_attachments(message)
            if attachment_text:
                content = f"{content}\n\n{attachment_text}".strip() if content else attachment_text

            image_attachments = await _download_image_attachments(message)
            audio_attachments = await _download_audio_attachments(message)
            video_attachments = await _download_video_attachments(message)

            if (
                not content.strip()
                and not image_attachments
                and not audio_attachments
                and not video_attachments
            ):
                return

            # Route through Awareness (blocking call - run in thread pool)
            async with message.channel.typing():
                packet = await asyncio.to_thread(
                    bot.awareness.synchronous_run,
                    content,
                    image_attachments=image_attachments,
                    audio_attachments=audio_attachments,
                    video_attachments=video_attachments,
                    text_attachments=text_attachments,
                )

            reply = bot.response_from_packet(packet)

            # Send reply in chunks if needed
            for index, chunk in enumerate(_split_reply(reply)):
                if index == 0:
                    await message.reply(chunk, mention_author=False)
                else:
                    await message.channel.send(chunk)

        except Exception as exc:
            print(f"[discord] Error processing message {message.id}: {exc}")
            await message.reply(
                "Hit an error processing that. Check the bot window for details.",
                mention_author=False,
            )


# ── Entry point ────────────────────────────────────────────────────────────


def main() -> None:
    if not DISCORD_TOKEN:
        print(
            "\n[discord] No token found.\n"
            "Add this line to world_models/config.py:\n\n"
            "    DISCORD_BOT_TOKEN = '<your bot token>'\n\n"
            "Or set the environment variable GOPHER_DISCORD_TOKEN."
        )
        raise SystemExit(1)

    print(f"[discord] Starting bridge (channel: #{DISCORD_CHANNEL})")
    client.run(DISCORD_TOKEN)


if __name__ == "__main__":
    main()
